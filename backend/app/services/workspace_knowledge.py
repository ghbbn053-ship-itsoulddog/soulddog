"""
工作区知识库服务：
- 工作区管理
- 用户上传文档入库
- 文档切块和向量化
- 轻量关系抽取（Skill/MCP/文档）
"""

from __future__ import annotations

import csv
import hashlib
import io
import json
import logging
import mimetypes
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from sqlalchemy.orm import Session

from app.models import User
from app.models.platform import (
    Workspace,
    KnowledgeSource,
    KnowledgeDocument,
    KnowledgeChunk,
    KnowledgeRelation,
)
from app.services import get_model_provider, get_vector_store

logger = logging.getLogger(__name__)

UPLOAD_ROOT = Path("backend/data/workspace_uploads")

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover - optional dependency
    PdfReader = None

try:
    from docx import Document as DocxDocument
except Exception:  # pragma: no cover - optional dependency
    DocxDocument = None

try:
    from openpyxl import load_workbook
except Exception:  # pragma: no cover - optional dependency
    load_workbook = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None


def _slugify(text: str) -> str:
    base = re.sub(r"[^a-zA-Z0-9\u4e00-\u9fff_-]+", "-", (text or "").strip()).strip("-").lower()
    return base or "workspace"


def _normalize_text(text: str) -> str:
    raw = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    raw = re.sub(r"\n{3,}", "\n\n", raw)
    lines = [line.strip() for line in raw.split("\n")]
    return "\n".join(lines).strip()


def _hash_bytes(raw_bytes: bytes) -> str:
    return hashlib.sha256(raw_bytes or b"").hexdigest()


def _split_text(text: str, chunk_size: int = 900, overlap: int = 120) -> List[str]:
    normalized = _normalize_text(text)
    if not normalized:
        return []
    if len(normalized) <= chunk_size:
        return [normalized]
    chunks: List[str] = []
    start = 0
    while start < len(normalized):
        end = min(len(normalized), start + chunk_size)
        chunk = normalized[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(normalized):
            break
        start = max(end - overlap, start + 1)
    return chunks


class WorkspaceKnowledgeService:
    def ensure_default_workspace(self, db: Session, owner_username: str) -> Workspace:
        workspace = (
            db.query(Workspace)
            .filter(Workspace.owner_username == owner_username, Workspace.is_default == True)
            .first()
        )
        if workspace:
            return workspace

        workspace = Workspace(
            owner_username=owner_username,
            slug="default",
            name="默认工作区",
            description="平台默认知识空间",
            is_default=True,
        )
        db.add(workspace)
        db.commit()
        db.refresh(workspace)
        return workspace

    def list_workspaces(self, db: Session, owner_username: str) -> List[Workspace]:
        self.ensure_default_workspace(db, owner_username)
        return (
            db.query(Workspace)
            .filter(Workspace.owner_username == owner_username)
            .order_by(Workspace.is_default.desc(), Workspace.id.asc())
            .all()
        )

    def create_workspace(self, db: Session, owner_username: str, name: str, description: str = "") -> Workspace:
        self.ensure_default_workspace(db, owner_username)
        slug_seed = _slugify(name)
        slug = slug_seed
        index = 2
        while (
            db.query(Workspace)
            .filter(Workspace.owner_username == owner_username, Workspace.slug == slug)
            .first()
        ):
            slug = f"{slug_seed}-{index}"
            index += 1
        workspace = Workspace(
            owner_username=owner_username,
            slug=slug,
            name=name.strip() or "未命名工作区",
            description=(description or "").strip(),
            is_default=False,
        )
        db.add(workspace)
        db.commit()
        db.refresh(workspace)
        return workspace

    def _owner_dir(self, owner_username: str, workspace_slug: str) -> Path:
        target = UPLOAD_ROOT / owner_username / workspace_slug
        target.mkdir(parents=True, exist_ok=True)
        return target

    def _extract_relations(self, owner_username: str, workspace_id: int, title: str, content: str) -> List[KnowledgeRelation]:
        relations: List[KnowledgeRelation] = []
        lowered = f"{title}\n{content}".lower()
        if "skill" in lowered:
            relations.append(
                KnowledgeRelation(
                    workspace_id=workspace_id,
                    owner_username=owner_username,
                    subject_type="document",
                    subject_ref=title,
                    relation_type="mentions",
                    object_type="capability",
                    object_ref="skill",
                    confidence="medium",
                    metadata_json={},
                )
            )
        if "mcp" in lowered:
            relations.append(
                KnowledgeRelation(
                    workspace_id=workspace_id,
                    owner_username=owner_username,
                    subject_type="document",
                    subject_ref=title,
                    relation_type="mentions",
                    object_type="capability",
                    object_ref="mcp",
                    confidence="medium",
                    metadata_json={},
                )
            )
        return relations

    @staticmethod
    def _decode_text_bytes(raw_bytes: bytes) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"):
            try:
                return raw_bytes.decode(encoding)
            except Exception:
                continue
        return raw_bytes.decode("utf-8", errors="ignore")

    def _extract_pdf_text(self, raw_bytes: bytes) -> str:
        if PdfReader is None:
            raise ValueError("当前未安装 PDF 解析依赖 pypdf")
        reader = PdfReader(io.BytesIO(raw_bytes))
        parts: List[str] = []
        for idx, page in enumerate(reader.pages):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                parts.append(f"[第 {idx + 1} 页]\n{page_text}")
        return "\n\n".join(parts).strip()

    def _extract_docx_text(self, raw_bytes: bytes) -> str:
        if DocxDocument is None:
            raise ValueError("当前未安装 DOCX 解析依赖 python-docx")
        doc = DocxDocument(io.BytesIO(raw_bytes))
        parts: List[str] = []
        for paragraph in doc.paragraphs:
            text = (paragraph.text or "").strip()
            if text:
                parts.append(text)
        for table_idx, table in enumerate(doc.tables):
            table_rows: List[str] = []
            for row in table.rows:
                cells = [(cell.text or "").strip() for cell in row.cells]
                if any(cells):
                    table_rows.append(" | ".join(cells))
            if table_rows:
                parts.append(f"[表格 {table_idx + 1}]\n" + "\n".join(table_rows))
        return "\n\n".join(parts).strip()

    def _extract_xlsx_text(self, raw_bytes: bytes) -> str:
        if load_workbook is None:
            raise ValueError("当前未安装 Excel 解析依赖 openpyxl")
        workbook = load_workbook(io.BytesIO(raw_bytes), data_only=True)
        parts: List[str] = []
        for sheet in workbook.worksheets:
            sheet_rows: List[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell).strip() for cell in row if cell not in (None, "")]
                if cells:
                    sheet_rows.append(" | ".join(cells))
            if sheet_rows:
                parts.append(f"[工作表] {sheet.title}\n" + "\n".join(sheet_rows))
        return "\n\n".join(parts).strip()

    def _extract_pptx_text(self, raw_bytes: bytes) -> str:
        if Presentation is None:
            raise ValueError("当前未安装 PPTX 解析依赖 python-pptx")
        prs = Presentation(io.BytesIO(raw_bytes))
        slides: List[str] = []
        for slide_idx, slide in enumerate(prs.slides):
            texts: List[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "") or ""
                text = text.strip()
                if text:
                    texts.append(text)
            if texts:
                slides.append(f"[幻灯片 {slide_idx + 1}]\n" + "\n".join(texts))
        return "\n\n".join(slides).strip()

    def _extract_csv_text(self, raw_bytes: bytes) -> str:
        text = self._decode_text_bytes(raw_bytes)
        reader = csv.reader(io.StringIO(text))
        rows: List[str] = []
        for row in reader:
            cells = [cell.strip() for cell in row if cell and cell.strip()]
            if cells:
                rows.append(" | ".join(cells))
        return "\n".join(rows).strip()

    def _extract_upload_text(
        self,
        filename: str,
        raw_bytes: bytes,
        mime_type: Optional[str] = None,
    ) -> Tuple[str, Dict]:
        suffix = Path(filename).suffix.lower()
        parser = "text"
        if suffix == ".json":
            parser = "json"
            parsed = json.loads(self._decode_text_bytes(raw_bytes))
            return json.dumps(parsed, ensure_ascii=False, indent=2), {"parser": parser}
        if suffix in {".txt", ".md", ".markdown", ".yaml", ".yml", ".html", ".htm", ".xml", ".log"}:
            return self._decode_text_bytes(raw_bytes), {"parser": parser}
        if suffix == ".csv" or mime_type == "text/csv":
            parser = "csv"
            return self._extract_csv_text(raw_bytes), {"parser": parser}
        if suffix == ".pdf":
            parser = "pdf"
            return self._extract_pdf_text(raw_bytes), {"parser": parser}
        if suffix == ".docx":
            parser = "docx"
            return self._extract_docx_text(raw_bytes), {"parser": parser}
        if suffix in {".xlsx", ".xlsm"}:
            parser = "xlsx"
            return self._extract_xlsx_text(raw_bytes), {"parser": parser}
        if suffix == ".pptx":
            parser = "pptx"
            return self._extract_pptx_text(raw_bytes), {"parser": parser}
        return self._decode_text_bytes(raw_bytes), {"parser": parser}

    def _list_source_documents(
        self,
        db: Session,
        owner_username: str,
        workspace_id: int,
        logical_name: str,
    ) -> List[KnowledgeDocument]:
        docs = (
            db.query(KnowledgeDocument)
            .join(KnowledgeSource, KnowledgeSource.id == KnowledgeDocument.source_id)
            .filter(
                KnowledgeDocument.owner_username == owner_username,
                KnowledgeDocument.workspace_id == workspace_id,
                KnowledgeSource.original_filename == logical_name,
            )
            .order_by(KnowledgeDocument.id.asc())
            .all()
        )
        return docs

    def _find_duplicate_document(
        self,
        db: Session,
        owner_username: str,
        workspace_id: int,
        content_hash: str,
    ) -> Optional[KnowledgeDocument]:
        docs = (
            db.query(KnowledgeDocument)
            .filter(
                KnowledgeDocument.owner_username == owner_username,
                KnowledgeDocument.workspace_id == workspace_id,
            )
            .order_by(KnowledgeDocument.id.desc())
            .all()
        )
        for doc in docs:
            meta = doc.metadata_json or {}
            if str(meta.get("content_hash", "")) == content_hash:
                return doc
        return None

    def _next_version(
        self,
        db: Session,
        owner_username: str,
        workspace_id: int,
        logical_name: str,
    ) -> int:
        docs = self._list_source_documents(db, owner_username, workspace_id, logical_name)
        version = 1
        for doc in docs:
            current = int((doc.metadata_json or {}).get("version", 1) or 1)
            if current >= version:
                version = current + 1
        return version

    def _persist_failed_document(
        self,
        db: Session,
        owner_username: str,
        workspace: Workspace,
        filename: str,
        raw_bytes: bytes,
        source_type: str,
        mime_type: Optional[str],
        authority_level: str,
        error_message: str,
        extra_meta: Optional[Dict] = None,
    ) -> KnowledgeDocument:
        version = self._next_version(db, owner_username, workspace.id, filename)
        owner_dir = self._owner_dir(owner_username, workspace.slug)
        storage_name = filename if version == 1 else f"{Path(filename).stem}.v{version}{Path(filename).suffix}"
        target_path = owner_dir / storage_name
        target_path.write_bytes(raw_bytes or b"")

        detected_mime = mime_type or mimetypes.guess_type(filename)[0] or "application/octet-stream"
        content_hash = _hash_bytes(raw_bytes or b"")
        source = KnowledgeSource(
            workspace_id=workspace.id,
            owner_username=owner_username,
            source_type=source_type,
            title=filename,
            mime_type=detected_mime,
            original_filename=filename,
            storage_path=str(target_path),
            status="failed",
            authority_level=authority_level,
            meta_json={
                "content_hash": content_hash,
                "logical_name": filename,
                "storage_name": storage_name,
                "version": version,
                "file_size": len(raw_bytes or b""),
                "error": error_message,
                **(extra_meta or {}),
            },
        )
        db.add(source)
        db.flush()

        doc = KnowledgeDocument(
            workspace_id=workspace.id,
            source_id=source.id,
            owner_username=owner_username,
            title=Path(filename).stem or filename,
            doc_type="document",
            status="failed",
            summary=f"文档解析失败：{error_message}",
            content_text="",
            token_estimate=0,
            metadata_json={
                "filename": filename,
                "source_type": source_type,
                "mime_type": detected_mime,
                "content_hash": content_hash,
                "version": version,
                "authority_level": authority_level,
                "storage_name": storage_name,
                "error": error_message,
                **(extra_meta or {}),
            },
        )
        db.add(doc)
        db.commit()
        db.refresh(doc)
        return doc

    def ingest_text_document(
        self,
        db: Session,
        owner_username: str,
        workspace_id: int,
        filename: str,
        content_text: str,
        source_type: str = "upload",
        mime_type: Optional[str] = None,
        extra_meta: Optional[Dict] = None,
        authority_level: str = "user",
    ) -> KnowledgeDocument:
        workspace = (
            db.query(Workspace)
            .filter(Workspace.id == workspace_id, Workspace.owner_username == owner_username)
            .first()
        )
        if not workspace:
            raise ValueError("工作区不存在")

        normalized = _normalize_text(content_text)
        if not normalized:
            raise ValueError("文档内容为空")

        raw_bytes = normalized.encode("utf-8")
        content_hash = _hash_bytes(raw_bytes)
        duplicate = self._find_duplicate_document(db, owner_username, workspace.id, content_hash)
        if duplicate:
            return duplicate

        version = self._next_version(db, owner_username, workspace.id, filename)
        owner_dir = self._owner_dir(owner_username, workspace.slug)
        storage_name = filename if version == 1 else f"{Path(filename).stem}.v{version}{Path(filename).suffix}"
        target_path = owner_dir / storage_name
        target_path.write_text(normalized, encoding="utf-8")

        detected_mime = mime_type or mimetypes.guess_type(filename)[0] or "text/plain"
        source = KnowledgeSource(
            workspace_id=workspace.id,
            owner_username=owner_username,
            source_type=source_type,
            title=filename,
            mime_type=detected_mime,
            original_filename=filename,
            storage_path=str(target_path),
            status="processing",
            authority_level=authority_level,
            meta_json={
                "content_hash": content_hash,
                "logical_name": filename,
                "storage_name": storage_name,
                "version": version,
                "file_size": len(raw_bytes),
                **(extra_meta or {}),
            },
        )
        db.add(source)
        db.flush()

        title = Path(filename).stem or filename
        summary = normalized[:240]
        doc = KnowledgeDocument(
            workspace_id=workspace.id,
            source_id=source.id,
            owner_username=owner_username,
            title=title,
            doc_type="document",
            status="ready",
            summary=summary,
            content_text=normalized,
            token_estimate=max(1, len(normalized) // 2),
            metadata_json={
                "filename": filename,
                "source_type": source_type,
                "mime_type": detected_mime,
                "content_hash": content_hash,
                "version": version,
                "authority_level": authority_level,
                "storage_name": storage_name,
                **(extra_meta or {}),
            },
        )
        db.add(doc)
        db.flush()

        chunks = _split_text(normalized)
        chunk_rows: List[KnowledgeChunk] = []
        for idx, chunk in enumerate(chunks):
            chunk_rows.append(
                KnowledgeChunk(
                    document_id=doc.id,
                    workspace_id=workspace.id,
                    owner_username=owner_username,
                    chunk_index=idx,
                    title=title if idx == 0 else f"{title} #{idx + 1}",
                    content=chunk,
                    char_count=len(chunk),
                    metadata_json={"document_title": title, "filename": filename},
                )
            )
        db.add_all(chunk_rows)

        relations = self._extract_relations(owner_username, workspace.id, title, normalized)
        if relations:
            db.add_all(relations)

        source.status = "ready"
        db.commit()
        db.refresh(doc)

        self._vectorize_document(doc, chunk_rows)
        return doc

    def ingest_uploaded_document(
        self,
        db: Session,
        owner_username: str,
        workspace_id: int,
        filename: str,
        raw_bytes: bytes,
        source_type: str = "upload",
        mime_type: Optional[str] = None,
        authority_level: str = "user",
        extra_meta: Optional[Dict] = None,
    ) -> KnowledgeDocument:
        if not raw_bytes:
            raise ValueError("上传文件为空")
        workspace = (
            db.query(Workspace)
            .filter(Workspace.id == workspace_id, Workspace.owner_username == owner_username)
            .first()
        )
        if not workspace:
            raise ValueError("工作区不存在")

        try:
            extracted_text, parse_meta = self._extract_upload_text(filename, raw_bytes, mime_type=mime_type)
            return self.ingest_text_document(
                db=db,
                owner_username=owner_username,
                workspace_id=workspace_id,
                filename=filename,
                content_text=extracted_text,
                source_type=source_type,
                mime_type=mime_type,
                authority_level=authority_level,
                extra_meta={
                    "original_file_size": len(raw_bytes),
                    **parse_meta,
                    **(extra_meta or {}),
                },
            )
        except ValueError as e:
            return self._persist_failed_document(
                db=db,
                owner_username=owner_username,
                workspace=workspace,
                filename=filename,
                raw_bytes=raw_bytes,
                source_type=source_type,
                mime_type=mime_type,
                authority_level=authority_level,
                error_message=str(e),
                extra_meta={"original_file_size": len(raw_bytes), **(extra_meta or {})},
            )

    def _vectorize_document(self, document: KnowledgeDocument, chunks: List[KnowledgeChunk]):
        if not chunks:
            return
        provider = get_model_provider()
        vec_store = get_vector_store()
        texts: List[str] = []
        embeddings: List[List[float]] = []
        sources: List[str] = []
        metadatas: List[Dict] = []

        for chunk in chunks:
            emb = provider.generate_embedding(chunk.content)
            if not emb:
                continue
            texts.append(chunk.content)
            embeddings.append(emb)
            sources.append(f"workspace:{document.workspace_id}:{document.title}")
            metadatas.append(
                {
                    "type": "knowledge_chunk",
                    "document_id": document.id,
                    "workspace_id": document.workspace_id,
                    "owner_username": document.owner_username,
                    "chunk_index": chunk.chunk_index,
                    "title": chunk.title,
                }
            )
        if texts and embeddings:
            user_numeric = abs(hash(document.owner_username)) % 2_000_000_000
            try:
                vec_store.add_documents(user_numeric, texts, embeddings, sources, metadatas)
            except Exception as e:
                logger.warning("工作区文档向量化失败: %s", e)

    def list_documents(self, db: Session, owner_username: str, workspace_id: Optional[int] = None) -> List[KnowledgeDocument]:
        query = db.query(KnowledgeDocument).filter(KnowledgeDocument.owner_username == owner_username)
        if workspace_id:
            query = query.filter(KnowledgeDocument.workspace_id == workspace_id)
        docs = query.order_by(KnowledgeDocument.id.desc()).all()
        hash_counts: Dict[str, int] = {}
        latest_version_by_name: Dict[str, int] = {}

        for doc in docs:
            meta = doc.metadata_json or {}
            content_hash = str(meta.get("content_hash", "")).strip()
            logical_name = str(meta.get("filename", "")).strip()
            version = int(meta.get("version", 1) or 1)
            if content_hash:
                hash_counts[content_hash] = hash_counts.get(content_hash, 0) + 1
            if logical_name:
                latest_version_by_name[logical_name] = max(latest_version_by_name.get(logical_name, 1), version)

        for doc in docs:
            meta = dict(doc.metadata_json or {})
            content_hash = str(meta.get("content_hash", "")).strip()
            logical_name = str(meta.get("filename", "")).strip()
            version = int(meta.get("version", 1) or 1)
            meta["duplicate_count"] = hash_counts.get(content_hash, 1) if content_hash else 1
            meta["is_duplicate"] = bool(content_hash and hash_counts.get(content_hash, 0) > 1)
            meta["latest_version"] = latest_version_by_name.get(logical_name, version) if logical_name else version
            meta["is_latest_version"] = meta["latest_version"] == version
            doc.metadata_json = meta
        return docs

    def get_workspace_knowledge_overview(self, db: Session, owner_username: str, workspace_id: int) -> Dict:
        workspace = (
            db.query(Workspace)
            .filter(Workspace.id == workspace_id, Workspace.owner_username == owner_username)
            .first()
        )
        if not workspace:
            raise ValueError("工作区不存在")

        docs = self.list_documents(db, owner_username, workspace_id=workspace_id)
        relation_count = (
            db.query(KnowledgeRelation)
            .filter(KnowledgeRelation.workspace_id == workspace_id, KnowledgeRelation.owner_username == owner_username)
            .count()
        )
        chunk_count = (
            db.query(KnowledgeChunk)
            .filter(KnowledgeChunk.workspace_id == workspace_id, KnowledgeChunk.owner_username == owner_username)
            .count()
        )

        by_doc_type: Dict[str, int] = {}
        by_status: Dict[str, int] = {}
        by_authority: Dict[str, int] = {}
        total_tokens = 0

        serialized_documents: List[Dict] = []
        for doc in docs:
            meta = dict(doc.metadata_json or {})
            authority_level = str(meta.get("authority_level", "user") or "user")
            by_doc_type[doc.doc_type] = by_doc_type.get(doc.doc_type, 0) + 1
            by_status[doc.status] = by_status.get(doc.status, 0) + 1
            by_authority[authority_level] = by_authority.get(authority_level, 0) + 1
            total_tokens += int(doc.token_estimate or 0)
            serialized_documents.append(
                {
                    "id": doc.id,
                    "workspace_id": doc.workspace_id,
                    "title": doc.title,
                    "doc_type": doc.doc_type,
                    "summary": doc.summary,
                    "status": doc.status,
                    "token_estimate": doc.token_estimate,
                    "metadata": meta,
                    "created_at": doc.created_at.isoformat() if doc.created_at else None,
                    "updated_at": doc.updated_at.isoformat() if doc.updated_at else None,
                }
            )

        return {
            "workspace": {
                "id": workspace.id,
                "slug": workspace.slug,
                "name": workspace.name,
                "description": workspace.description,
                "is_default": workspace.is_default,
            },
            "stats": {
                "documents": len(docs),
                "knowledge_units": chunk_count,
                "relations": relation_count,
                "ready_documents": by_status.get("ready", 0),
                "failed_documents": by_status.get("failed", 0),
                "total_tokens": total_tokens,
                "by_doc_type": by_doc_type,
                "by_status": by_status,
                "by_authority": by_authority,
            },
            "documents": serialized_documents,
        }

    def list_document_chunks(self, db: Session, owner_username: str, workspace_id: int, document_id: int) -> List[Dict]:
        workspace = (
            db.query(Workspace)
            .filter(Workspace.id == workspace_id, Workspace.owner_username == owner_username)
            .first()
        )
        if not workspace:
            raise ValueError("工作区不存在")

        document = (
            db.query(KnowledgeDocument)
            .filter(
                KnowledgeDocument.id == document_id,
                KnowledgeDocument.workspace_id == workspace_id,
                KnowledgeDocument.owner_username == owner_username,
            )
            .first()
        )
        if not document:
            raise ValueError("文档不存在")

        chunks = (
            db.query(KnowledgeChunk)
            .filter(
                KnowledgeChunk.document_id == document_id,
                KnowledgeChunk.workspace_id == workspace_id,
                KnowledgeChunk.owner_username == owner_username,
            )
            .order_by(KnowledgeChunk.chunk_index.asc())
            .all()
        )
        return [
            {
                "id": chunk.id,
                "document_id": document_id,
                "workspace_id": workspace_id,
                "chunk_index": chunk.chunk_index,
                "title": chunk.title,
                "content": chunk.content,
                "char_count": chunk.char_count,
                "metadata": chunk.metadata_json or {},
            }
            for chunk in chunks
        ]

    def delete_document(self, db: Session, owner_username: str, workspace_id: int, document_id: int) -> Dict:
        workspace = (
            db.query(Workspace)
            .filter(Workspace.id == workspace_id, Workspace.owner_username == owner_username)
            .first()
        )
        if not workspace:
            raise ValueError("工作区不存在")

        document = (
            db.query(KnowledgeDocument)
            .filter(
                KnowledgeDocument.id == document_id,
                KnowledgeDocument.workspace_id == workspace_id,
                KnowledgeDocument.owner_username == owner_username,
            )
            .first()
        )
        if not document:
            raise ValueError("文档不存在")

        source = (
            db.query(KnowledgeSource)
            .filter(
                KnowledgeSource.id == document.source_id,
                KnowledgeSource.workspace_id == workspace_id,
                KnowledgeSource.owner_username == owner_username,
            )
            .first()
        )
        title = document.title
        storage_path = source.storage_path if source else None

        db.delete(document)
        if source:
            remaining = (
                db.query(KnowledgeDocument)
                .filter(KnowledgeDocument.source_id == source.id)
                .count()
            )
            if remaining <= 0:
                db.delete(source)
        db.commit()

        try:
            if storage_path:
                path = Path(storage_path)
                if path.exists():
                    path.unlink()
        except Exception as e:
            logger.warning("删除工作区文档文件失败: %s", e)

        return {"success": True, "title": title}

    @staticmethod
    def _fallback_search_score(query: str, text: str) -> float:
        q = (query or "").strip().lower()
        t = (text or "").strip().lower()
        if not q or not t:
            return 0.0
        if q in t:
            return min(1.0, 0.65 + (len(q) / max(len(t), 1)))
        q_terms = [term for term in re.split(r"[\s,，。；;、]+", q) if term]
        if not q_terms:
            return 0.0
        hits = sum(1 for term in q_terms if term in t)
        return hits / max(len(q_terms), 1)

    def search_workspace(
        self,
        db: Session,
        owner_username: str,
        workspace_id: int,
        query: str,
        top_k: int = 8,
    ) -> List[Dict]:
        workspace = (
            db.query(Workspace)
            .filter(Workspace.id == workspace_id, Workspace.owner_username == owner_username)
            .first()
        )
        if not workspace:
            raise ValueError("工作区不存在")

        provider = get_model_provider()
        vec_store = get_vector_store()
        query_embedding = provider.generate_embedding(query)
        user_numeric = abs(hash(owner_username)) % 2_000_000_000
        hits: List[Dict] = []

        if query_embedding:
            try:
                raw_hits = vec_store.search(user_numeric, query_embedding, top_k=top_k * 2)
                for hit in raw_hits:
                    meta = hit.get("metadata") or {}
                    if int(meta.get("workspace_id", -1)) != int(workspace_id):
                        continue
                    hits.append(
                        {
                            "document_id": meta.get("document_id"),
                            "title": meta.get("title") or hit.get("source") or "未知文档",
                            "text": hit.get("text", ""),
                            "content": hit.get("text", ""),
                            "score": float(hit.get("score", 0.0)),
                            "chunk_index": meta.get("chunk_index", 0),
                            "source": hit.get("source", ""),
                        }
                    )
                    if len(hits) >= top_k:
                        break
            except Exception as e:
                logger.warning("工作区向量检索失败，回退文本匹配: %s", e)

        if hits:
            return hits[:top_k]

        docs = (
            db.query(KnowledgeDocument)
            .filter(KnowledgeDocument.owner_username == owner_username, KnowledgeDocument.workspace_id == workspace_id)
            .all()
        )
        scored: List[Tuple[float, KnowledgeDocument]] = []
        for doc in docs:
            score = self._fallback_search_score(query, f"{doc.title}\n{doc.summary or ''}\n{doc.content_text or ''}")
            if score > 0:
                scored.append((score, doc))
        scored.sort(key=lambda item: item[0], reverse=True)
        return [
            {
                "document_id": doc.id,
                "title": doc.title,
                "text": (doc.summary or doc.content_text or "")[:400],
                "content": (doc.summary or doc.content_text or "")[:400],
                "score": float(score),
                "chunk_index": 0,
                "source": f"workspace:{workspace_id}:{doc.title}",
            }
            for score, doc in scored[:top_k]
        ]

    def get_workspace_graph(self, db: Session, owner_username: str, workspace_id: int) -> Dict:
        workspace = (
            db.query(Workspace)
            .filter(Workspace.id == workspace_id, Workspace.owner_username == owner_username)
            .first()
        )
        if not workspace:
            raise ValueError("工作区不存在")
        docs = (
            db.query(KnowledgeDocument)
            .filter(KnowledgeDocument.workspace_id == workspace_id, KnowledgeDocument.owner_username == owner_username)
            .all()
        )
        rels = (
            db.query(KnowledgeRelation)
            .filter(KnowledgeRelation.workspace_id == workspace_id, KnowledgeRelation.owner_username == owner_username)
            .order_by(KnowledgeRelation.id.desc())
            .all()
        )
        nodes = [
            {"id": f"doc:{doc.id}", "label": doc.title, "type": doc.doc_type}
            for doc in docs
        ]
        node_ids = {node["id"] for node in nodes}
        edge_ids = set()
        edges = [
            {
                "id": f"rel:{rel.id}",
                "source": rel.subject_ref if rel.subject_ref.startswith("doc:") else f"label:{rel.subject_ref}",
                "target": rel.object_ref if rel.object_ref.startswith("doc:") else f"label:{rel.object_ref}",
                "label": rel.relation_type,
                "subject_type": rel.subject_type,
                "object_type": rel.object_type,
            }
            for rel in rels
        ]
        edge_ids.update(edge["id"] for edge in edges)

        for edge in edges:
            if edge["source"] not in node_ids:
                nodes.append({"id": edge["source"], "label": edge["source"].split(":", 1)[-1], "type": "label"})
                node_ids.add(edge["source"])
            if edge["target"] not in node_ids:
                nodes.append({"id": edge["target"], "label": edge["target"].split(":", 1)[-1], "type": "label"})
                node_ids.add(edge["target"])

        from app.services.skill_manager import get_skill_manager
        from app.services.mcp_registry import get_mcp_registry
        from app.services.composition_manager import get_composition_manager

        skill_items = get_skill_manager().list_skills(owner_username)
        mcp_items = get_mcp_registry().list_tools()
        composition = get_composition_manager().resolved(owner_username)

        for skill in skill_items:
            node_id = f"skill:{skill['name']}"
            if node_id not in node_ids:
                nodes.append({"id": node_id, "label": skill["name"], "type": "skill"})
                node_ids.add(node_id)
            edge_id = f"derived:workspace-skill:{workspace_id}:{skill['name']}"
            if edge_id not in edge_ids:
                edges.append(
                    {
                        "id": edge_id,
                        "source": f"workspace:{workspace_id}",
                        "target": node_id,
                        "label": "contains",
                        "subject_type": "workspace",
                        "object_type": "skill",
                    }
                )
                edge_ids.add(edge_id)

        for tool in mcp_items:
            tool_name = str(tool.get("name", "")).strip()
            if not tool_name:
                continue
            node_id = f"mcp:{tool_name}"
            if node_id not in node_ids:
                nodes.append({"id": node_id, "label": tool_name, "type": "mcp"})
                node_ids.add(node_id)
            edge_id = f"derived:workspace-mcp:{workspace_id}:{tool_name}"
            if edge_id not in edge_ids:
                edges.append(
                    {
                        "id": edge_id,
                        "source": f"workspace:{workspace_id}",
                        "target": node_id,
                        "label": "can_use",
                        "subject_type": "workspace",
                        "object_type": "mcp",
                    }
                )
                edge_ids.add(edge_id)

        if f"workspace:{workspace_id}" not in node_ids:
            nodes.append({"id": f"workspace:{workspace_id}", "label": workspace.name, "type": "workspace"})
            node_ids.add(f"workspace:{workspace_id}")

        for skill_name in composition.get("skills", []):
            skill_node = f"skill:{skill_name}"
            if skill_node in node_ids:
                edge_id = f"derived:enabled-skill:{workspace_id}:{skill_name}"
                if edge_id not in edge_ids:
                    edges.append(
                        {
                            "id": edge_id,
                            "source": f"workspace:{workspace_id}",
                            "target": skill_node,
                            "label": "enabled",
                            "subject_type": "workspace",
                            "object_type": "skill",
                        }
                    )
                    edge_ids.add(edge_id)

        for tool_name in composition.get("mcp_tools", []):
            tool_node = f"mcp:{tool_name}"
            if tool_node in node_ids:
                edge_id = f"derived:enabled-mcp:{workspace_id}:{tool_name}"
                if edge_id not in edge_ids:
                    edges.append(
                        {
                            "id": edge_id,
                            "source": f"workspace:{workspace_id}",
                            "target": tool_node,
                            "label": "enabled",
                            "subject_type": "workspace",
                            "object_type": "mcp",
                        }
                    )
                    edge_ids.add(edge_id)
        return {
            "workspace": {"id": workspace.id, "name": workspace.name, "slug": workspace.slug},
            "nodes": nodes,
            "edges": edges,
        }


_workspace_knowledge_singleton: Optional[WorkspaceKnowledgeService] = None


def get_workspace_knowledge_service() -> WorkspaceKnowledgeService:
    global _workspace_knowledge_singleton
    if _workspace_knowledge_singleton is None:
        _workspace_knowledge_singleton = WorkspaceKnowledgeService()
    return _workspace_knowledge_singleton
